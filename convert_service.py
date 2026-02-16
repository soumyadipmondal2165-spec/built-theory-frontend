import img2pdf
from PIL import Image
import io
from services.helpers import unique_path
from pdf2docx import Converter
from pptx import Presentation

# jpg/png images -> pdf
def images_to_pdf(file_streams):
    # file_streams: list of file-like objects
    images_bytes = [fs.read() for fs in file_streams]
    outpath = unique_path(".pdf")
    with open(outpath, "wb") as f:
        f.write(img2pdf.convert(images_bytes))
    return outpath

# pdf -> jpg (first page)
import fitz
def pdf_to_jpg(file_stream):
    doc = fitz.open(stream=file_stream.read(), filetype="pdf")
    page = doc.load_page(0)
    pix = page.get_pixmap()
    outpath = unique_path(".jpg")
    pix.save(outpath)
    doc.close()
    return outpath

# pdf -> docx
def pdf_to_word(file_stream):
    temp_in = unique_path(".pdf")
    with open(temp_in, "wb") as f:
        f.write(file_stream.read())
    outdoc = unique_path(".docx")
    cv = Converter(temp_in)
    cv.convert(outdoc)
    cv.close()
    return outdoc

# simple ppt generator
def generate_pptx(title, text):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = text or ""
    outpath = unique_path(".pptx")
    prs.save(outpath)
    return outpath
