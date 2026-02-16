import os
import pandas as pd
from docx2pdf import convert as docx_to_pdf_convert
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from services.helpers import unique_path

# 1. Word to PDF
def word_to_pdf(file_stream):
    # docx2pdf usually requires an actual file on disk
    temp_in = unique_path(".docx")
    temp_out = unique_path(".pdf")
    
    with open(temp_in, "wb") as f:
        f.write(file_stream.read())
    
    # Note: On Linux servers, this may require 'LibreOffice' installed
    docx_to_pdf_convert(temp_in, temp_out)
    return temp_out

# 2. Excel to PDF
def excel_to_pdf(file_stream):
    # Read Excel data
    df = pd.read_excel(file_stream)
    outpath = unique_path(".pdf")
    
    # Create a simple PDF table representation
    c = canvas.Canvas(outpath, pagesize=letter)
    width, height = letter
    y = height - 40
    
    c.setFont("Helvetica-Bold", 12)
    c.drawString(30, y, "Excel Data Export (Built-Theory)")
    y -= 30
    
    c.setFont("Helvetica", 10)
    # Simple row-by-row print (for advanced tables, use Platypus)
    for index, row in df.head(20).iterrows(): # Limits to 20 rows for performance
        row_str = " | ".join([str(val) for val in row.values])
        c.drawString(30, y, row_str[:100]) # Truncate long lines
        y -= 15
        if y < 40:
            c.showPage()
            y = height - 40
            
    c.save()
    return outpath

# 3. PPT to PDF (Text Extraction)
def ppt_to_pdf(file_stream):
    prs = Presentation(file_stream)
    outpath = unique_path(".pdf")
    c = canvas.Canvas(outpath, pagesize=letter)
    width, height = letter
    
    for i, slide in enumerate(prs.slides):
        y = height - 50
        c.setFont("Helvetica-Bold", 14)
        c.drawString(30, y, f"Slide {i+1}")
        y -= 30
        
        c.setFont("Helvetica", 11)
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                c.drawString(30, y, shape.text[:90])
                y -= 20
        c.showPage()
        
    c.save()
    return outpath